"""
agent.py — Auto-PPT Agent (Enhanced: LangChain + Title/Filename Fixes)

Flow:
  User Input
  → LangChain message objects (HumanMessage / SystemMessage)
  → clean_topic() + safe_filename()
  → Agentic ReAct loop (Plan → Create → AddSlides → WriteText → Save)
  → Post-process: apply dark-blue theme + white text styling
  → .pptx saved to generated_ppts/

Enhancements (over original working version):
  1. LangChain message schema (HumanMessage / SystemMessage) used to build
     and manage conversation messages — replaces raw dict construction.
     Groq still drives the LLM; LangChain is used only for message handling.
  2. Stronger clean_topic() with extra filler patterns so the topic is ALWAYS clean.
  3. Hardened title injection: topic + filename are bolted into the system
     prompt AND the user message so the LLM cannot ignore them.
  4. Runtime guard: write_text calls for slide_index=0 are intercepted and
     the title is forcibly corrected before reaching the MCP server.
  5. post_process_presentation() forcibly overwrites slide-0 title text
     after the file is written — last-resort guarantee against prompt leakage.

Rules kept intact:
  - ppt_mcp_server.py is NOT modified.
  - MCP integration (HTTP transport, session, tool calls) is unchanged.
  - Groq LLM and agentic loop logic are unchanged.
  - app.py import (run_ppt_agent) continues to work with no changes.
"""

import os
import re
import sys
import json
import asyncio

from groq import Groq
from mcp import ClientSession
from mcp.client.streamable_http import streamablehttp_client

# ── LangChain message schema (minimal, safe integration) ──────────────────────
# LangChain is used ONLY for structured message objects (HumanMessage /
# SystemMessage).  The actual LLM call is still handled by the Groq client.
# This satisfies the LangChain requirement without rewriting any agent logic.
from langchain_core.messages import HumanMessage, SystemMessage

# python-pptx styling imports
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# ──────────────────────────────────────────────
#  CLIENT SETUP
# ──────────────────────────────────────────────
client         = Groq(api_key=os.getenv("GROQ_API_KEY"))
MODEL_NAME     = "llama-3.3-70b-versatile"
MAX_ITERATIONS = 30          # hard cap — prevents infinite loops


# ──────────────────────────────────────────────
#  HELPERS
# ──────────────────────────────────────────────

def mcp_tool_to_groq_format(mcp_tool):
    """Convert an MCP tool definition into the Groq function-call schema."""
    return {
        "type": "function",
        "function": {
            "name":        mcp_tool.name,
            "description": mcp_tool.description,
            "parameters":  mcp_tool.inputSchema,
        },
    }


def clean_topic(user_request: str) -> str:
    """
    Strip common filler phrases from the user request to obtain a clean topic name.

    Enhanced with additional patterns to handle more edge cases.

    Examples
    --------
    "create a ppt of 6 slides on artificial intelligence"  → "Artificial Intelligence"
    "make a 5-slide deck on Solar System"                  → "Solar System"
    "generate presentation about quantum computing"        → "Quantum Computing"
    "create a ppt on AI"                                   → "Ai"  (title-cased)
    """
    filler_patterns = [
        # "create/make/generate [a] [N-slide] ppt/presentation [of N slides] on/about/for"
        r"(?:please\s+)?(?:create|make|generate|build|design|produce)\s+"
        r"(?:a\s+)?(?:(?:\d+[\s\-]*(?:slides?|slide)[\s\-]*)?)"
        r"(?:ppt|powerpoint|presentation|deck|slides?)\s+"
        r"(?:of\s+\d+\s+slides?\s+)?(?:on|about|for|regarding)\s+",

        # "create/make/generate [slides] on/about/for"
        r"(?:create|make|generate|build|design)\s+(?:slides?\s+)?(?:on|about|for)\s+",

        # "a ppt on", "a deck on", "a presentation on"
        r"a\s+(?:ppt|deck|presentation)\s+(?:on|about|for)\s+",

        # "presentation/ppt/slides on/about/for"
        r"(?:presentation|ppt|slides?)\s+(?:on|about|for|regarding)\s+",

        # trailing slide-count fragments: "of 6 slides", "with 5 slides"
        r"\s+(?:of|with)\s+\d+\s+slides?\s*$",

        # leftover leading articles after all other stripping
        r"^(?:a|an|the)\s+",
    ]

    cleaned = user_request.strip()
    for pattern in filler_patterns:
        cleaned = re.sub(pattern, "", cleaned, flags=re.IGNORECASE).strip()

    # Remove stray leading/trailing punctuation
    cleaned = cleaned.strip(" .,;:-")

    # Fallback: if cleaning wiped everything, use original title-cased input
    result = cleaned.title() if cleaned else user_request.strip().title()
    return result


def safe_filename(topic: str) -> str:
    """
    Convert a topic string into a safe .pptx filename.

    Rules:
      - lowercase
      - spaces / hyphens → underscores
      - non-word characters removed
      - suffix: _presentation.pptx

    Example: "Artificial Intelligence" → "artificial_intelligence_presentation.pptx"
    """
    name = re.sub(r"[^\w\s-]", "", topic).strip().lower()
    name = re.sub(r"[\s-]+", "_", name)
    return f"{name}_presentation.pptx"


# ──────────────────────────────────────────────
#  LANGCHAIN MESSAGE BUILDERS
# ──────────────────────────────────────────────

def build_system_message(topic: str, out_name: str) -> SystemMessage:
    """
    Build a LangChain SystemMessage that embeds the full agent instructions.

    The clean topic and filename are injected directly into the prompt text
    so the LLM has no ambiguity about what to write or where to save.

    Returns a LangChain SystemMessage object.
    """
    content = f"""
You are an expert Presentation Agent. Your only job is to create a complete,
high-quality PowerPoint presentation using the tools available to you.

═══════════════════════════════════════════════════════
MANDATORY EXECUTION ORDER — NEVER DEVIATE
═══════════════════════════════════════════════════════

STEP 1 ── PLAN (output as plain text BEFORE any tool call)
  Write exactly this block first:

  PLAN:
  - Title Slide : {topic}
  - Slide 1     : [2–4 word title]
  - Slide 2     : [2–4 word title]
  - Slide 3     : [2–4 word title]
  - Slide 4     : [2–4 word title]
  - Slide 5     : [2–4 word title]

STEP 2 ── Call create_presentation()

STEP 3 ── TITLE SLIDE  ⚠️ STRICT RULES BELOW
  Call add_slide()
  Call write_text(
      slide_index = 0,
      title       = "{topic}",
      bullets     = []
  )
  ⚠️ title MUST be exactly: "{topic}"
  ⚠️ bullets MUST be [] — NO bullets on the title slide EVER

STEP 4 ── CONTENT SLIDES  (repeat for every planned slide)
  Call add_slide()
  Call write_text(
      slide_index = N,
      title       = "[2–4 word title]",
      bullets     = ["bullet 1", "bullet 2", "bullet 3", "bullet 4"]
  )

STEP 5 ── SAVE
  Call save_presentation(filename="{out_name}")

═══════════════════════════════════════════════════════
TITLE SLIDE RULES
═══════════════════════════════════════════════════════
• title for slide_index 0 MUST be exactly: "{topic}"
• bullets MUST be [] for slide_index 0.
• NEVER write the raw user prompt as the title.

  ✔  title = "{topic}",  bullets = []
  ✗  title = "Create a ppt on {topic}"       ← WRONG
  ✗  title = "Presentation about {topic}"    ← WRONG

═══════════════════════════════════════════════════════
SLIDE TITLE RULES  (content slides)
═══════════════════════════════════════════════════════
• 2 to 4 words maximum.
• Specific and meaningful — use title case.

═══════════════════════════════════════════════════════
BULLET POINT RULES
═══════════════════════════════════════════════════════
• Exactly 3 to 4 bullets per content slide.
• Each bullet: 8 to 15 words.
• Format: [Subject] + [what it does / means / why it matters].
• Language: simple, clear, professional — no jargon.
• NO repetition across slides. Stay strictly on topic.

═══════════════════════════════════════════════════════
FILENAME RULE
═══════════════════════════════════════════════════════
• Save filename MUST be exactly: "{out_name}"

═══════════════════════════════════════════════════════
ERROR HANDLING
═══════════════════════════════════════════════════════
• On tool error: read the message, fix arguments, retry.
• NEVER stop because of a single failure.
• ALWAYS end with save_presentation — no exceptions.
"""
    return SystemMessage(content=content)


def build_user_message(topic: str, out_name: str, user_request: str) -> HumanMessage:
    """
    Build a LangChain HumanMessage that states the clean topic and filename
    explicitly, so the LLM cannot misread the original user prompt as a title.

    Returns a LangChain HumanMessage object.
    """
    content = (
        f"Create a complete PowerPoint presentation on: '{topic}'\n\n"
        f"User's original request (for context only): '{user_request}'\n\n"
        f"⚠️  TITLE SLIDE (slide_index=0) title MUST be exactly: \"{topic}\"\n"
        f"⚠️  Title slide bullets MUST be an empty list: []\n"
        f"⚠️  Save filename MUST be exactly: \"{out_name}\""
    )
    return HumanMessage(content=content)


def langchain_message_to_groq_dict(msg) -> dict:
    """
    Convert a LangChain BaseMessage (SystemMessage / HumanMessage) into the
    plain dict format expected by the Groq chat completions API.

    This is the ONLY bridge between LangChain and Groq — zero other LangChain
    APIs are invoked, so the existing agentic loop is completely unchanged.
    """
    if isinstance(msg, SystemMessage):
        return {"role": "system", "content": msg.content}
    if isinstance(msg, HumanMessage):
        return {"role": "user",   "content": msg.content}
    # Generic fallback
    return {"role": "user", "content": str(msg.content)}


# ──────────────────────────────────────────────
#  STYLING POST-PROCESSOR
# ──────────────────────────────────────────────

_BG_COLOR   = RGBColor(0x03, 0x25, 0x6C)   # #03256C dark blue
_TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)    # white
_TITLE_PT   = Pt(35)
_BULLET_PT  = Pt(15)


def _set_slide_background(slide, color: RGBColor) -> None:
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _style_text_frame(tf, font_size: Pt, bold: bool = False) -> None:
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size      = font_size
            run.font.bold      = bold
            run.font.color.rgb = _TEXT_WHITE
        para.font.size      = font_size
        para.font.bold      = bold
        para.font.color.rgb = _TEXT_WHITE


def post_process_presentation(filepath: str, topic: str = "") -> None:
    """
    Open the saved .pptx, apply the brand theme to every slide, and re-save.

    Enhancement: if `topic` is provided, slide-0 title text is FORCIBLY
    overwritten with the clean topic string — last-resort guard against any
    prompt text leaking through from the LLM.

    Theme applied:
      • Background  : #03256C (dark blue)
      • Title text  : white, 35 pt, bold
      • Bullet text : white, 15 pt
      • Slide 0     : title forced to `topic`; content placeholder cleared
    """
    if not os.path.isfile(filepath):
        print(f"[Style] File not found, skipping: {filepath}")
        return

    try:
        prs = Presentation(filepath)

        for idx, slide in enumerate(prs.slides):

            _set_slide_background(slide, _BG_COLOR)

            title_ph = slide.shapes.title
            if title_ph and title_ph.has_text_frame:

                # ✅ Forcibly overwrite slide-0 title with clean topic
                if idx == 0 and topic:
                    title_ph.text_frame.clear()
                    title_ph.text_frame.text = topic

                _style_text_frame(title_ph.text_frame, _TITLE_PT, bold=True)

            content_ph = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    content_ph = shape
                    break

            if idx == 0:
                if content_ph and content_ph.has_text_frame:
                    content_ph.text_frame.clear()
                    content_ph.text_frame.text = ""
            else:
                if content_ph and content_ph.has_text_frame:
                    _style_text_frame(content_ph.text_frame, _BULLET_PT, bold=False)

        prs.save(filepath)
        print(f"[Style] ✅ Styling applied → {filepath}")

    except Exception as style_err:
        print(
            f"[Style] ⚠️  Post-processing failed "
            f"(presentation still saved): {style_err}"
        )


# ──────────────────────────────────────────────
#  MAIN AGENT  (async)
# ──────────────────────────────────────────────

async def run_ppt_agent(user_request: str) -> str:
    """
    Run the full agentic loop and return a status string.

    LangChain integration summary:
      - build_system_message() returns a LangChain SystemMessage.
      - build_user_message()   returns a LangChain HumanMessage.
      - langchain_message_to_groq_dict() converts both to Groq-compatible dicts.
      - All subsequent messages in the loop are plain dicts (unchanged from original).
      - The Groq client, MCP session, and tool execution loop are untouched.

    MCP connection (unchanged):
      Connects to the already-running ppt_mcp_server.py over HTTP:
          http://127.0.0.1:8000/mcp
      Start the server first:
          python ppt_mcp_server.py
    """

    # ── 1. Extract clean topic + filename ─────────────────────────────────
    topic    = clean_topic(user_request)
    out_name = safe_filename(topic)

    print(f"\n{'='*55}")
    print(f"  Auto-PPT Agent  (LangChain + MCP + Groq)")
    print(f"  Topic    : {topic}")
    print(f"  Output   : generated_ppts/{out_name}")
    print(f"{'='*55}\n")

    # ── 2. Build LangChain message objects ────────────────────────────────
    system_lc_msg = build_system_message(topic, out_name)
    user_lc_msg   = build_user_message(topic, out_name, user_request)

    print(f"[LangChain] SystemMessage built  ✓  ({len(system_lc_msg.content)} chars)")
    print(f"[LangChain] HumanMessage  built  ✓  ({len(user_lc_msg.content)} chars)")

    # ── 3. Convert LangChain messages → Groq dicts ────────────────────────
    messages = [
        langchain_message_to_groq_dict(system_lc_msg),
        langchain_message_to_groq_dict(user_lc_msg),
    ]

    # ── 4. MCP HTTP connection ─────────────────────────────────────────────
    MCP_URL = "http://127.0.0.1:8000/mcp"
    print(f"[*] Connecting to MCP server at {MCP_URL} …")

    async with streamablehttp_client(MCP_URL) as (read, write, _):
        async with ClientSession(read, write) as session:

            await session.initialize()

            tools_response = await session.list_tools()
            groq_tools     = [mcp_tool_to_groq_format(t) for t in tools_response.tools]
            tool_names     = [t.name for t in tools_response.tools]

            print(f"[*] MCP connected  ✓   Tools: {tool_names}\n")
            print(f"[*] Agentic loop started …\n")

            # ── 5. Agentic ReAct loop ──────────────────────────────────────
            iteration   = 0
            save_called = False

            while iteration < MAX_ITERATIONS and not save_called:
                iteration += 1
                print(f"┌─ Iteration {iteration}/{MAX_ITERATIONS} {'─'*40}")

                try:
                    response = client.chat.completions.create(
                        model=MODEL_NAME,
                        messages=messages,
                        tools=groq_tools,
                        tool_choice="auto",
                        max_tokens=4096,
                    )
                except Exception as llm_err:
                    print(f"│  [LLM Error] {llm_err}")
                    print(f"└─ Retrying next iteration …\n")
                    continue

                msg      = response.choices[0].message
                msg_dict = {"role": msg.role}

                if msg.content:
                    msg_dict["content"] = msg.content
                    snippet = msg.content[:400].replace("\n", " ")
                    print(f"│  [Thought] {snippet}{'…' if len(msg.content) > 400 else ''}")

                if msg.tool_calls:
                    tool_calls_dict = [
                        {
                            "id":   tc.id,
                            "type": "function",
                            "function": {
                                "name":      tc.function.name,
                                "arguments": tc.function.arguments,
                            },
                        }
                        for tc in msg.tool_calls
                    ]
                    msg_dict["tool_calls"] = tool_calls_dict
                    messages.append(msg_dict)

                    for tool_call in msg.tool_calls:
                        func_name = tool_call.function.name

                        try:
                            func_args = json.loads(tool_call.function.arguments)
                        except json.JSONDecodeError:
                            func_args = {}

                        # ── Runtime guard: intercept title-slide write ─────
                        # If the LLM passes the wrong title for slide 0,
                        # silently correct it before calling the MCP server.
                        if func_name == "write_text":
                            if func_args.get("slide_index") == 0:
                                current_title = func_args.get("title", "").strip()
                                if current_title.lower() != topic.lower():
                                    print(
                                        f"│  [Guard]  Title corrected: "
                                        f'"{current_title}" → "{topic}"'
                                    )
                                    func_args["title"]   = topic
                                    func_args["bullets"] = []

                        # ── Intercept save_presentation ───────────────────
                        if func_name == "save_presentation":
                            os.makedirs("generated_ppts", exist_ok=True)
                            raw_fn   = func_args.get("filename", out_name)
                            clean_fn = os.path.basename(raw_fn)
                            if not clean_fn.endswith(".pptx"):
                                clean_fn += ".pptx"
                            func_args["filename"] = os.path.join(
                                "generated_ppts", clean_fn
                            )

                        print(f"│  [Action]      {func_name}")
                        arg_preview = json.dumps(func_args)
                        if len(arg_preview) > 160:
                            arg_preview = arg_preview[:160] + "…"
                        print(f"│  [Args]        {arg_preview}")

                        try:
                            result      = await session.call_tool(func_name, func_args)
                            result_text = (
                                result.content[0].text
                                if result.content
                                else "Success (no output)"
                            )
                        except Exception as tool_err:
                            result_text = f"Tool error: {tool_err}"

                        print(f"│  [Observation] {result_text}")

                        messages.append({
                            "role":         "tool",
                            "tool_call_id": tool_call.id,
                            "content":      result_text,
                        })

                        if (
                            func_name == "save_presentation"
                            and "saved at" in result_text.lower()
                        ):
                            save_called = True

                            saved_path = func_args.get("filename", "")
                            abs_match  = re.search(
                                r"saved at:\s*(.+\.pptx)",
                                result_text,
                                re.IGNORECASE,
                            )
                            if abs_match:
                                saved_path = abs_match.group(1).strip()

                            if saved_path and os.path.isfile(saved_path):
                                post_process_presentation(saved_path, topic=topic)
                            elif saved_path:
                                alt = os.path.abspath(saved_path)
                                if os.path.isfile(alt):
                                    post_process_presentation(alt, topic=topic)

                            print(f"└─\n")
                            print(f"{'='*55}")
                            print(f"  ✅ SUCCESS — Presentation saved & styled!")
                            print(f"  Path: {func_args['filename']}")
                            print(f"{'='*55}\n")
                            return result_text

                else:
                    messages.append({"role": "assistant", "content": msg.content or ""})

                    already_saved = any(
                        m.get("role") == "tool"
                        and "saved at" in m.get("content", "").lower()
                        for m in messages
                    )

                    if already_saved:
                        save_called = True
                        print(f"└─\n✅ [SUCCESS] Detected prior save — finishing.\n")
                        return "Presentation saved successfully."

                    nudge = (
                        "You must continue using the tools provided. "
                        "Do NOT narrate — ACT. "
                    )
                    called_create = any(
                        m.get("role") == "tool"
                        and "created successfully" in m.get("content", "").lower()
                        for m in messages
                    )
                    if not called_create:
                        nudge += "Call create_presentation() right now."
                    else:
                        nudge += (
                            "Add the next slide with add_slide(), fill it with "
                            "write_text(), and once ALL slides are done, call "
                            "save_presentation() to finish."
                        )

                    print(f"│  [System] {nudge}")
                    messages.append({"role": "user", "content": nudge})

                print(f"└─\n")

            # ── FALLBACK SAVE ──────────────────────────────────────────────
            if not save_called:
                print(
                    "[FALLBACK] Loop ended without confirmed save "
                    "— attempting emergency save …"
                )
                os.makedirs("generated_ppts", exist_ok=True)
                fallback_path = os.path.join("generated_ppts", out_name)
                try:
                    result   = await session.call_tool(
                        "save_presentation", {"filename": fallback_path}
                    )
                    fb_text  = result.content[0].text if result.content else "Done"
                    print(f"[FALLBACK] {fb_text}")
                    if "saved at" in fb_text.lower():
                        abs_m = re.search(
                            r"saved at:\s*(.+\.pptx)", fb_text, re.IGNORECASE
                        )
                        if abs_m:
                            post_process_presentation(abs_m.group(1).strip(), topic=topic)
                        elif os.path.isfile(fallback_path):
                            post_process_presentation(fallback_path, topic=topic)
                    return fb_text
                except Exception as fb_err:
                    msg_out = f"Fallback save failed: {fb_err}"
                    print(f"[FALLBACK ERROR] {msg_out}")
                    return msg_out

    return "Agent loop completed."


# ──────────────────────────────────────────────
#  CLI ENTRY POINT
# ──────────────────────────────────────────────

if __name__ == "__main__":
    if len(sys.argv) > 1:
        req = " ".join(sys.argv[1:])
    else:
        try:
            req = input("Enter your PPT topic: ").strip()
        except EOFError:
            req = "The lifecycle of a star for a 6th-grade class"

    if not req:
        req = "Introduction to Artificial Intelligence"

    asyncio.run(run_ppt_agent(req))
