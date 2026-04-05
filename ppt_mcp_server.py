"""
ppt_mcp_server.py

MCP Server for PowerPoint Automation using python-pptx.

Provides tools to:
- Create a presentation
- Add slides
- Write text content
- Insert images
- Save presentation

Designed for use with an AI Agent (Auto-PPT Agent).
"""

from fastmcp import FastMCP
from pptx import Presentation
from pptx.util import Inches
import requests
from io import BytesIO
import os


# Initialize MCP server
mcp = FastMCP("ppt-server")

# Global presentation object (maintains state across tool calls)
presentation = {"prs": None}


# -------------------------------
# 1. CREATE PRESENTATION
# -------------------------------
@mcp.tool()
def create_presentation():
    """
    Initialize a new empty PowerPoint presentation.

    Must be called before adding slides.

    Returns:
        str: Status message
    """
    presentation["prs"] = Presentation()
    return "Presentation created successfully"


# -------------------------------
# 2. ADD SLIDE
# -------------------------------
@mcp.tool()
def add_slide():
    """
    Add a new slide using 'Title + Content' layout.

    Returns:
        int | str: Slide index if successful, else error message
    """
    prs = presentation["prs"]

    if prs is None:
        return "Error: Create presentation first"

    layout = prs.slide_layouts[1]  # Title + Content
    slide = prs.slides.add_slide(layout)

    return len(prs.slides) - 1


# -------------------------------
# 3. WRITE TEXT
# -------------------------------
@mcp.tool()
def write_text(slide_index: int, title: str, bullets: list):
    """
    Add title and bullet points to a specific slide.

    Args:
        slide_index (int): Index of the slide
        title (str): Slide title
        bullets (list): List of bullet point strings

    Returns:
        str: Status message
    """
    prs = presentation["prs"]

    if prs is None:
        return "Error: No presentation found"

    if slide_index >= len(prs.slides):
        return "Error: Invalid slide index"

    slide = prs.slides[slide_index]

    # Set title
    slide.shapes.title.text = title

    # Set bullet content
    tf = slide.placeholders[1].text_frame
    tf.clear()

    if not bullets:
        return "Warning: No bullets provided"

    # First bullet (important fix)
    tf.text = bullets[0]

    # Remaining bullets
    for bullet in bullets[1:]:
        p = tf.add_paragraph()
        p.text = bullet

    return f"Text added to slide {slide_index}"


# -------------------------------
# 4. ADD IMAGE
# -------------------------------
@mcp.tool()
def add_image(slide_index: int, image_url: str):
    """
    Add an image to a slide from a URL.

    Args:
        slide_index (int): Index of the slide
        image_url (str): Direct image URL

    Returns:
        str: Status message
    """
    prs = presentation["prs"]

    if prs is None:
        return "Error: No presentation found"

    if slide_index >= len(prs.slides):
        return "Error: Invalid slide index"

    slide = prs.slides[slide_index]

    try:
        response = requests.get(image_url, timeout=5)

        if response.status_code != 200:
            return "Error: Invalid image URL"

        image_stream = BytesIO(response.content)

        # Positioning image (right side of slide)
        slide.shapes.add_picture(
            image_stream,
            Inches(5),
            Inches(2),
            width=Inches(4)
        )

        return "Image added successfully"

    except Exception:
        return "Error: Failed to add image"


# -------------------------------
# 5. SAVE PRESENTATION
# -------------------------------
@mcp.tool()
def save_presentation(filename: str = "output.pptx"):
    """
    Save the presentation to local disk.

    Args:
        filename (str): Output file name

    Returns:
        str: File save confirmation
    """
    prs = presentation["prs"]

    if prs is None:
        return "Error: Nothing to save"

    try:
        path = os.path.join(os.getcwd(), filename)
        prs.save(path)
        return f"Presentation saved at: {path}"

    except Exception:
        return "Error: Failed to save presentation"


# -------------------------------
# 6. GET SLIDE COUNT (Optional Helper)
# -------------------------------
@mcp.tool()
def get_slide_count():
    """
    Get the current number of slides.

    Returns:
        int | str: Number of slides or error message
    """
    prs = presentation["prs"]

    if prs is None:
        return "Error: No presentation found"

    return len(prs.slides)


# -------------------------------
# RUN MCP SERVER
# -------------------------------
if __name__ == "__main__":
    mcp.run(transport="http", port=8000)